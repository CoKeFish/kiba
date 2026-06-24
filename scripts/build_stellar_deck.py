# -*- coding: utf-8 -*-
"""Genera el pitch deck de Kiba planteado sobre Stellar (.pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Paleta de marca ----
NAVY      = RGBColor(0x02, 0x05, 0x0E)   # fondo profundo
NAVY2     = RGBColor(0x05, 0x0B, 0x1B)   # fondo alterno
PANEL     = RGBColor(0x1B, 0x27, 0x50)   # tarjetas
PANEL_LT  = RGBColor(0x2A, 0x38, 0x68)
BLUE      = RGBColor(0x7A, 0xA3, 0xFF)   # acento claro  (landing --blue-300)
BLUE_DK   = RGBColor(0x20, 0x60, 0xF6)   # acento primario (landing --accent)
GREEN     = RGBColor(0x2E, 0xD3, 0x9A)   # acento secundario
GOLD      = RGBColor(0xF5, 0xB5, 0x44)
PINK      = RGBColor(0xFF, 0x55, 0x71)
WHITE     = RGBColor(0xF7, 0xF8, 0xFB)
MUTED     = RGBColor(0xA9, 0xB4, 0xCC)
MUTED_DK  = RGBColor(0x7B, 0x89, 0xA8)

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
LOGO = os.path.join(ROOT, "packages", "landing", "public", "logomark.png")
OUT  = os.path.join(ROOT, "docs", "Kiba-Stellar.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.06):
    """runs: lista de párrafos; cada párrafo es lista de (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Segoe UI"
    return tb


def kicker(s, txt, color=BLUE, x=0.9, y=0.62, sq=None):
    box(s, x, y + 0.02, 0.32, 0.32, fill=(sq if sq is not None else color), radius=False)
    text(s, x + 0.45, y, 8, 0.4, [[(txt.upper(), 13, color, True)]])


def title(s, txt, x=0.9, y=1.0, w=11.5, size=33):
    text(s, x, y, w, 1.1, [[(txt, size, WHITE, True)]])


def footer(s, idx):
    text(s, 0.9, 7.02, 6, 0.3, [[("kiba  ·  on Stellar", 10, MUTED_DK, False)]])
    text(s, 11.4, 7.02, 1.05, 0.3, [[(f"{idx:02d}", 10, MUTED_DK, True)]],
         align=PP_ALIGN.RIGHT)


def chip(s, x, y, txt, color):
    w = 0.30 + 0.092 * len(txt)
    box(s, x, y, w, 0.42, fill=PANEL, line=color, line_w=1.25)
    text(s, x, y + 0.015, w, 0.40, [[(txt, 11, color, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


def card(s, x, y, w, h, head, body, accent):
    box(s, x, y, w, h, fill=PANEL, line=PANEL_LT, line_w=1.0)
    box(s, x, y, 0.09, h, fill=accent, radius=False)
    text(s, x + 0.32, y + 0.26, w - 0.55, 0.5, [[(head, 16, WHITE, True)]])
    text(s, x + 0.32, y + 0.78, w - 0.55, h - 1.0,
         [[(body, 12.5, MUTED, False)]], line_spacing=1.12)


# ============================================================== 1 · PORTADA
s = slide(NAVY)
# acento diagonal
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.85), SW, Inches(0.65))
band.fill.solid(); band.fill.fore_color.rgb = BLUE_DK; band.line.fill.background()
band.shadow.inherit = False
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.9), Inches(0.85), height=Inches(1.25))
text(s, 0.9, 2.55, 11.5, 1.4, [[("kiba", 60, WHITE, True)]])
text(s, 0.92, 3.75, 11.5, 0.8,
     [[("Un marketplace donde los asistentes de IA descubren y pagan ", 21, MUTED, False)]])
text(s, 0.92, 4.18, 11.5, 0.8,
     [[("a agentes especializados bajo demanda — liquidado en ", 21, MUTED, False),
       ("Stellar", 21, GREEN, True), (".", 21, MUTED, False)]])
chip(s, 0.92, 5.25, "Soroban", GREEN)
chip(s, 2.15, 5.25, "USDC nativo", BLUE)
chip(s, 3.78, 5.25, "x402", GOLD)
chip(s, 4.92, 5.25, "Agentic commerce", PINK)
text(s, 0.9, 6.95, 11.5, 0.45,
     [[("Pitch de concepto  ·  Bogotá, Colombia  ·  2026", 12, WHITE, True)]])

# ============================================================== 2 · PROBLEMA
s = slide(NAVY2); kicker(s, "El problema", PINK)
title(s, "Los asistentes generales fallan en lo especializado")
text(s, 0.9, 2.05, 7.2, 2.4,
     [[("Un asistente de uso general (Claude, Cursor, ChatGPT) es bueno en tareas amplias,", 15, MUTED, False)],
      [("pero ", 15, MUTED, False), ("inventa o se equivoca", 15, WHITE, True),
       (" cuando la tarea exige experiencia concreta y datos en vivo.", 15, MUTED, False)],
      [("", 6, MUTED, False)],
      [("Conectar un servicio especializado hoy obliga a cada usuario a registrarse, leer docs,", 15, MUTED, False)],
      [("gestionar credenciales y escribir integraciones. ", 15, MUTED, False),
       ("Casi nadie lo hace.", 15, WHITE, True)]],
     line_spacing=1.15, space_after=4)
card(s, 8.35, 1.95, 4.1, 1.45, "El usuario",
     "Recibe respuestas confiadas pero erróneas en temas que importan.", PINK)
card(s, 8.35, 3.55, 4.1, 1.45, "El especialista",
     "Tiene la respuesta de calidad, pero no llega a esos usuarios.", GOLD)
text(s, 0.9, 5.55, 11.5, 0.9,
     [[("Resultado: los asistentes se quedan genéricos, los expertos invisibles, y la ", 16, MUTED, False),
       ("capacidad útil nunca encuentra a quien la necesita.", 16, WHITE, True)]],
     line_spacing=1.2)
footer(s, 2)

# ============================================================== 3 · SOLUCIÓN
s = slide(NAVY); kicker(s, "La solución", GREEN)
title(s, "Un único punto de entrada al ecosistema de agentes")
text(s, 0.9, 1.95, 11.5, 0.9,
     [[("El asistente localiza al agente experto, recibe un precio, paga en un solo paso y devuelve", 15, MUTED, False)],
      [("la respuesta — todo transparente para la persona. Sin claves, sin registros, sin billeteras.", 15, MUTED, False)]],
     line_spacing=1.15)
cards = [
    ("Descubrir", "El asistente busca por intención y encuentra al agente adecuado para la tarea.", BLUE),
    ("Pagar", "Micropago en USDC sobre Stellar, en la misma ida y vuelta de la consulta.", GREEN),
    ("Responder", "El experto entrega; el usuario recibe una respuesta confiable y citable.", GOLD),
]
for i, (h, b, a) in enumerate(cards):
    card(s, 0.9 + i * 4.0, 3.25, 3.7, 2.2, h, b, a)
text(s, 0.9, 5.85, 11.5, 0.8,
     [[("Dos lados, un protocolo: ", 15, WHITE, True),
       ("los usuarios acceden a expertos sin fricción; los publishers cobran sin construir infraestructura.", 15, MUTED, False)]],
     line_spacing=1.2)
footer(s, 3)

# ============================================================== 4 · CÓMO FUNCIONA
s = slide(NAVY2); kicker(s, "Cómo funciona", BLUE, sq=BLUE_DK)
title(s, "Del prompt al pago en una sola interacción")
steps = [
    ("1", "Intención", "El usuario pide algo especializado en su asistente de siempre."),
    ("2", "Match", "El marketplace encuentra al agente experto y devuelve un precio."),
    ("3", "Pago", "Se abre un escrow en Stellar y se paga en USDC al confirmar."),
    ("4", "Entrega", "El agente responde y cobra; el reparto se liquida en cadena."),
]
x = 0.9; w = 2.85; gap = 0.32
for i, (n, h, b) in enumerate(steps):
    cx = x + i * (w + gap)
    box(s, cx, 2.5, w, 2.6, fill=PANEL, line=PANEL_LT, line_w=1.0)
    box(s, cx + 0.28, 2.78, 0.62, 0.62, fill=BLUE_DK)
    text(s, cx + 0.28, 2.80, 0.62, 0.62, [[(n, 22, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 0.28, 3.62, w - 0.5, 0.5, [[(h, 16, WHITE, True)]])
    text(s, cx + 0.28, 4.10, w - 0.5, 0.9, [[(b, 12, MUTED, False)]], line_spacing=1.12)
    if i < 3:
        text(s, cx + w - 0.02, 3.55, 0.4, 0.5, [[("→", 22, BLUE, True)]],
             align=PP_ALIGN.CENTER)
text(s, 0.9, 5.55, 11.5, 0.7,
     [[("El protocolo de pago x402 (HTTP-nativo) maneja el ", 14, MUTED, False),
       ("402 Payment Required → quote → pago → entrega", 14, GREEN, True),
       (" sin sacar al usuario de la conversación.", 14, MUTED, False)]],
     line_spacing=1.2)
footer(s, 4)

# ============================================================== 5 · POR QUÉ STELLAR
s = slide(NAVY); kicker(s, "Por qué Stellar", GREEN)
title(s, "La red hecha para pagos es la red para agentes")
items = [
    ("USDC nativo", "Stablecoin de primera clase: micropagos sin volatilidad, contabilidad en dólares.", GREEN),
    ("Fees mínimas", "Fracciones de centavo por operación — viable cobrar por llamada, no por suscripción.", BLUE),
    ("Finalidad rápida", "Liquidación en ~5 segundos: el pago no interrumpe la conversación.", GOLD),
    ("Anchors / fiat", "Rampa de entrada y salida a moneda local para usuarios y publishers.", BLUE),
    ("Soroban", "Contratos en Rust para el escrow y el reparto de ingresos atómico.", PINK),
    ("Agentic commerce", "Foco estratégico de Stellar (y de Jed McCaleb) en comercio entre agentes.", GREEN),
]
gx, gy, gw, gh = 0.9, 2.1, 3.83, 1.5
for i, (h, b, a) in enumerate(items):
    cx = gx + (i % 3) * (gw + 0.22)
    cy = gy + (i // 3) * (gh + 0.22)
    card(s, cx, cy, gw, gh, h, b, a)
footer(s, 5)

# ============================================================== 6 · ARQUITECTURA
s = slide(NAVY2); kicker(s, "Arquitectura", BLUE, sq=BLUE_DK)
title(s, "Tres capas, un contrato que reparte solo")
layers = [
    ("Clientes", "Instalador de 1 clic o npx — MCP en Claude Desktop, Cursor y Claude Code; SDK y dashboard web.", BLUE),
    ("Plataforma", "Descubrimiento por intención, billeteras custodiadas y créditos en USDC.", GREEN),
    ("Stellar / Soroban", "Registro de agentes y escrow; el contrato libera el pago al entregar.", GOLD),
]
for i, (h, b, a) in enumerate(layers):
    cy = 2.0 + i * 1.18
    box(s, 0.9, cy, 7.4, 1.02, fill=PANEL, line=PANEL_LT, line_w=1.0)
    box(s, 0.9, cy, 0.09, 1.02, fill=a, radius=False)
    text(s, 1.2, cy + 0.13, 2.6, 0.8, [[(h, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, cy + 0.13, 4.5, 0.8, [[(b, 12, MUTED, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
# panel reparto
box(s, 8.65, 2.0, 3.78, 3.38, fill=PANEL, line=GREEN, line_w=1.25)
text(s, 8.95, 2.25, 3.3, 0.5, [[("Reparto atómico", 16, WHITE, True)]])
text(s, 8.95, 2.85, 3.3, 1.6,
     [[("95%", 40, GREEN, True), ("  al publisher", 15, MUTED, False)],
      [("5%", 40, BLUE, True), ("  a la plataforma", 15, MUTED, False)]],
     line_spacing=1.05, space_after=10)
text(s, 8.95, 4.55, 3.3, 0.8,
     [[("Aplicado dentro del contrato. Nadie puede desviar ingresos fuera de cadena.", 11.5, MUTED, False)]],
     line_spacing=1.12)
footer(s, 6)

# ============================================================== 7 · CÓMO SE CONECTA
s = slide(NAVY); kicker(s, "Cómo se conecta", BLUE, sq=BLUE_DK)
title(s, "De cero a conectado sin tocar la terminal")
text(s, 0.9, 1.95, 11.5, 0.6,
     [[("Tres caminos al mismo marketplace; la autenticación es por ", 15, MUTED, False),
       ("OAuth en el navegador", 15, WHITE, True),
       (" — sin pegar claves ni gestionar billeteras.", 15, MUTED, False)]],
     line_spacing=1.15)
ways = [
    ("Instalador de escritorio", "1 CLIC", GREEN,
     "Un .exe que detecta Claude Desktop, Cursor y Claude Code, respalda tu configuración e instala el MCP por ti.",
     None),
    ("npx / npm", "1 LÍNEA", BLUE,
     "Añade el servidor a tu config de MCP. Para quien ya vive en la terminal.",
     "npx -y kiba-mcp"),
    ("SDK", "PUBLISHERS", GOLD,
     "Integra tu propio agente y exponlo en el marketplace para cobrar por llamada.",
     None),
]
pw, py0, ph = 3.7, 2.65, 2.75
for i, (head, badge, accent, body, code) in enumerate(ways):
    px = 0.9 + i * 4.0
    box(s, px, py0, pw, ph, fill=PANEL, line=PANEL_LT, line_w=1.0)
    box(s, px, py0, 0.09, ph, fill=accent, radius=False)
    text(s, px + 0.32, py0 + 0.24, pw - 0.6, 0.45, [[(head, 15.5, WHITE, True)]])
    bw = 0.26 + 0.085 * len(badge)
    box(s, px + 0.32, py0 + 0.74, bw, 0.32, fill=None, line=accent, line_w=1.0)
    text(s, px + 0.32, py0 + 0.755, bw, 0.30, [[(badge, 9.5, accent, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, px + 0.32, py0 + 1.2, pw - 0.6, 1.1, [[(body, 12, MUTED, False)]],
         line_spacing=1.12)
    if code:
        box(s, px + 0.32, py0 + ph - 0.62, pw - 0.64, 0.42, fill=NAVY, line=accent, line_w=1.0)
        cb = s.shapes.add_textbox(Inches(px + 0.42), Inches(py0 + ph - 0.605),
                                  Inches(pw - 0.82), Inches(0.40))
        ctf = cb.text_frame; ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.LEFT
        cr = cp.add_run(); cr.text = code
        cr.font.size = Pt(11); cr.font.color.rgb = GREEN; cr.font.bold = True
        cr.font.name = "Consolas"
box(s, 0.9, 5.72, 11.55, 0.82, fill=PANEL, line=PANEL_LT, line_w=1.0)
text(s, 1.3, 5.86, 11.0, 0.6,
     [[("Hecho esto, el asistente recibe sus herramientas: ", 13, MUTED, False),
       ("list_agents · call_agent · get_balance · get_transactions", 13, BLUE, True)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
footer(s, 7)

# ============================================================== 8 · MERCADO DE DOS LADOS
s = slide(NAVY); kicker(s, "Mercado de dos lados", GOLD)
title(s, "Cada lado gana algo que hoy no tiene")
box(s, 0.9, 2.1, 5.65, 3.3, fill=PANEL, line=BLUE, line_w=1.25)
text(s, 1.25, 2.4, 5, 0.5, [[("Usuarios de asistentes", 18, BLUE, True)]])
for i, t in enumerate([
        "Capacidades expertas dentro de su chat de siempre.",
        "Sin registros por servicio, claves ni billeteras.",
        "Pagan solo por lo que usan, en dólares.",
        "Respuestas confiables y citables."]):
    text(s, 1.25, 3.0 + i * 0.56, 5.1, 0.5,
         [[("•  ", 14, BLUE, True), (t, 13.5, MUTED, False)]], line_spacing=1.1)
box(s, 6.8, 2.1, 5.65, 3.3, fill=PANEL, line=GREEN, line_w=1.25)
text(s, 7.15, 2.4, 5, 0.5, [[("Publishers de agentes", 18, GREEN, True)]])
for i, t in enumerate([
        "Distribución instantánea a todos los asistentes.",
        "Cobro y reparto resueltos por el protocolo.",
        "Sin construir facturación ni infraestructura.",
        "Monetizan su experiencia desde el primer día."]):
    text(s, 7.15, 3.0 + i * 0.56, 5.1, 0.5,
         [[("•  ", 14, GREEN, True), (t, 13.5, MUTED, False)]], line_spacing=1.1)
text(s, 0.9, 5.7, 11.5, 0.7,
     [[("El efecto de red: ", 15, WHITE, True),
       ("más agentes hacen al marketplace más útil, y más usuarios lo hacen más atractivo para publicar.", 15, MUTED, False)]],
     line_spacing=1.2)
footer(s, 8)

# ============================================================== 9 · ESTADO
s = slide(NAVY2); kicker(s, "Estado", GREEN)
title(s, "Qué existe hoy y qué sigue")
box(s, 0.9, 2.05, 5.65, 3.85, fill=PANEL, line=GREEN, line_w=1.0)
text(s, 1.25, 2.30, 5, 0.5, [[("Demostrado", 17, GREEN, True)]])
for i, t in enumerate([
        "Marketplace funcional de punta a punta.",
        "Descubrimiento por intención (palabra clave + semántico).",
        "Pago por llamada vía x402 con escrow on-chain.",
        "Reparto atómico aplicado por el contrato.",
        "Acceso por MCP en Claude Desktop, Cursor y Claude Code.",
        "Instalador de escritorio de 1 clic + paquete npm publicado."]):
    text(s, 1.25, 2.84 + i * 0.5, 5.1, 0.5,
         [[("✓  ", 13, GREEN, True), (t, 12.5, MUTED, False)]], line_spacing=1.08)
box(s, 6.8, 2.05, 5.65, 3.85, fill=PANEL, line=GOLD, line_w=1.0)
text(s, 7.15, 2.30, 5, 0.5, [[("Lo que sigue", 17, GOLD, True)]])
for i, t in enumerate([
        "Onboarding de publishers externos.",
        "Rampa fiat con anchors de Stellar.",
        "Salida a mainnet y auditoría del contrato Soroban.",
        "Catálogo de agentes verticales.",
        "Primeros usuarios pagos reales."]):
    text(s, 7.15, 2.84 + i * 0.5, 5.1, 0.5,
         [[("→  ", 13, GOLD, True), (t, 12.5, MUTED, False)]], line_spacing=1.08)
footer(s, 9)

# ============================================================== 10 · EQUIPO + CIERRE
s = slide(NAVY); kicker(s, "Equipo y visión", BLUE, sq=BLUE_DK)
title(s, "Construido en Bogotá, pensado para cualquiera")
team = [
    ("Rodion Tabares", "Ingeniería · plataforma, billeteras, descubrimiento, MCP"),
    ("André Landinez", "Ingeniería · contrato, pricing dinámico, traza de pago"),
    ("Lizeth Rico", "Diseño · identidad visual y experiencia de producto"),
]
for i, (n, r) in enumerate(team):
    cx = 0.9 + i * 4.0
    box(s, cx, 2.15, 3.7, 1.55, fill=PANEL, line=PANEL_LT, line_w=1.0)
    box(s, cx + 0.3, 2.45, 0.55, 0.55, fill=BLUE_DK)
    text(s, cx + 0.3, 3.12, 3.2, 0.4, [[(n, 14.5, WHITE, True)]])
    text(s, cx + 0.3, 3.5, 3.2, 0.5, [[(r, 10.5, MUTED, False)]], line_spacing=1.05)
box(s, 0.9, 4.1, 11.55, 2.05, fill=PANEL, line=GREEN, line_w=1.25)
text(s, 1.3, 4.45, 10.8, 0.6,
     [[("Es el momento de invertir energía en una herramienta que conecte a los asistentes de IA", 17, WHITE, True)]])
text(s, 1.3, 4.92, 10.8, 0.6,
     [[("con agentes y servicios especializados de forma efectiva.", 17, WHITE, True)]])
text(s, 1.3, 5.5, 10.8, 0.6,
     [[("Stellar nos da los pagos; nosotros ponemos el mercado.", 15, GREEN, True)]])
footer(s, 10)

prs.save(OUT)
print("OK:", OUT, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

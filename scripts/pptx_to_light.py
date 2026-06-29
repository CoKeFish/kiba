"""Remap the dark deck to a light theme by recoloring fills + font colors."""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor

SRC = "docs/Kiba-Stellar.pptx"
DST = "docs/Kiba-Stellar.pptx"

# Map: original hex (uppercase) -> new hex
FILL_MAP = {
    "02050E": "FFFFFF",  # main dark bg -> white
    "050B1B": "F4F6FB",  # alt dark bg  -> very light gray
    "1B2750": "EEF2FA",  # card surface -> light card
    # accent fills stay (kept vivid against white)
    "2060F6": "2060F6",
    "2ED39A": "1FB585",
    "F5B544": "E0A024",
    "FF5571": "E63E5C",
    "7AA3FF": "4E7CE8",
}

TEXT_MAP = {
    "F7F8FB": "0B1530",  # light text -> near-black
    "A9B4CC": "47536F",  # secondary  -> mid-dark gray
    "7B89A8": "6B7A9C",  # tertiary footer -> mid gray
    # accents darkened for AA contrast on white
    "7AA3FF": "1F4FBE",
    "2ED39A": "0E8A60",
    "F5B544": "8A5E0E",
    "FF5571": "C42E48",
    "2060F6": "1A4FCC",
}


def remap_fill(shape):
    try:
        fill = shape.fill
        if fill.type is None:
            return
        fc = fill.fore_color
        rgb = fc.rgb
        if rgb is None:
            return
        key = str(rgb).upper()
        if key in FILL_MAP:
            fc.rgb = RGBColor.from_string(FILL_MAP[key])
    except Exception:
        pass


def remap_runs(tf):
    for para in tf.paragraphs:
        for r in para.runs:
            try:
                rgb = r.font.color.rgb
                if rgb is None:
                    continue
                key = str(rgb).upper()
                if key in TEXT_MAP:
                    r.font.color.rgb = RGBColor.from_string(TEXT_MAP[key])
            except Exception:
                pass


def walk(shape):
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            walk(sub)
        return
    remap_fill(shape)
    if shape.has_text_frame:
        remap_runs(shape.text_frame)


p = Presentation(SRC)
for s in p.slides:
    for sh in s.shapes:
        walk(sh)

p.save(DST)
print("saved", DST)
